#!/usr/bin/env python3
"""Small, reproducible bridge between Figure Studio and PPT Master.

Every external command writes its complete stdout/stderr to ``quality/*.log``.
The public functions raise :class:`BridgeError` on failure so orchestrators can
name the failing stage without losing the underlying diagnostic output.
"""

from __future__ import annotations
import os

import argparse
import json
import re
import shutil
import subprocess
import sys
import uuid
from pathlib import Path
from typing import Any, Iterable
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE


PY = Path(os.environ.get("FIGURE_STUDIO_PYTHON", sys.executable))
PM = Path(os.environ.get("PPT_MASTER_HOME", Path.home() / ".argus-skill/tools/ppt-master/skills/ppt-master"))
PROJECT_MANAGER = PM / "scripts/project_manager.py"
SVG_CHECKER = PM / "scripts/svg_quality_checker.py"
SVG_TO_PPTX = PM / "scripts/svg_to_pptx.py"
PPTX_TO_SVG = PM / "scripts/pptx_to_svg.py"


class BridgeError(RuntimeError):
    """A PPT Master subprocess failed; ``log_path`` contains its transcript."""

    def __init__(self, message: str, *, command: Iterable[object], log_path: Path):
        super().__init__(message)
        self.command = [str(item) for item in command]
        self.log_path = log_path


def _quality_dir(path: Path) -> Path:
    """Resolve the nearest project quality directory for an input artifact."""
    resolved = path.resolve()
    if resolved.parent.name in {"svg_output", "roundtrip"}:
        project = resolved.parent.parent
    elif resolved.is_dir() and resolved.name == "roundtrip":
        project = resolved.parent
    else:
        project = resolved.parent
    quality = project / "quality"
    quality.mkdir(parents=True, exist_ok=True)
    return quality


def _run(
    command: Iterable[object],
    log_path: Path,
    *,
    cwd: Path | None = None,
    check: bool = True,
) -> subprocess.CompletedProcess[str]:
    argv = [str(item) for item in command]
    completed = subprocess.run(
        argv,
        cwd=str(cwd) if cwd else None,
        text=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=False,
    )
    log_path.parent.mkdir(parents=True, exist_ok=True)
    transcript = (
        "$ " + " ".join(argv) + "\n"
        + f"exit_code: {completed.returncode}\n"
        + "\n[stdout]\n" + completed.stdout
        + "\n[stderr]\n" + completed.stderr
    )
    log_path.write_text(transcript, encoding="utf-8")
    if check and completed.returncode != 0:
        raise BridgeError(
            f"command failed with exit code {completed.returncode}",
            command=argv,
            log_path=log_path,
        )
    return completed


_HEX_COLOR = re.compile(r"#[0-9A-Fa-f]{6}\b")
_COLOR_ROLES = {
    "#FBFAF7": "background",
    "#FFFFFF": "background",
    "#1F2933": "primary",
    "#D55E00": "accent",
    "#111827": "text",
    "#4B5563": "text_secondary",
    "#9CA3AF": "group_border",
    "#CBD5E1": "muted_line",
    "#D9DEE5": "pale_line",
    "#FFE2D1": "input",
    "#FFF2BD": "process",
    "#DCECFF": "memory",
    "#E2F7DF": "agent",
    "#EADFFF": "output",
    "#FFF1C9": "benchmark",
    "#F3F4F6": "neutral",
    "#FFF3EC": "input_group",
    "#FFF9DF": "process_group",
    "#EEF6FF": "memory_group",
    "#F0FAEE": "agent_group",
    "#F5F0FF": "output_group",
    "#FFF8E5": "benchmark_group",
    "#F8F9FA": "neutral_group",
}
_COLOR_ROLE_ORDER = {role: index for index, role in enumerate(_COLOR_ROLES.values())}


def _number(value: float) -> str:
    return str(int(value)) if abs(value - round(value)) < 1e-9 else f"{value:.2f}".rstrip("0").rstrip(".")


def _svg_inventory(svg_path: Path) -> dict[str, Any]:
    """Collect the exact palette, typography, and Tabler icons emitted by SVG."""

    root = ET.parse(svg_path).getroot()
    colors: set[str] = set()
    font_families: set[str] = set()
    font_sizes: set[float] = set()
    sizes_by_role: dict[str, set[float]] = {}
    icons: set[str] = set()
    paint_attributes = {"fill", "stroke", "stop-color", "flood-color", "color"}
    for element in root.iter():
        for attribute in paint_attributes:
            colors.update(match.group(0).upper() for match in _HEX_COLOR.finditer(element.get(attribute, "")))
        style = element.get("style", "")
        colors.update(match.group(0).upper() for match in _HEX_COLOR.finditer(style))
        if element.get("font-family"):
            font_families.add(str(element.get("font-family")).strip())
        if element.get("font-size"):
            size = float(str(element.get("font-size")))
            font_sizes.add(size)
            role = element.get("data-text-role") or element.get("data-figure-role") or "text"
            sizes_by_role.setdefault(str(role), set()).add(size)
        href = element.get("href") or element.get("{http://www.w3.org/1999/xlink}href") or ""
        if href.startswith("#ic-"):
            icons.add(href[4:])
    if not colors or not font_families or not font_sizes:
        raise ValueError(f"rendered SVG inventory is incomplete: {svg_path}")
    return {
        "colors": sorted(
            colors,
            key=lambda color: (_COLOR_ROLE_ORDER.get(_COLOR_ROLES.get(color, ""), len(_COLOR_ROLE_ORDER)), color),
        ),
        "font_families": sorted(font_families),
        "font_sizes": sorted(font_sizes),
        "sizes_by_role": {role: sorted(sizes) for role, sizes in sorted(sizes_by_role.items())},
        "icons": sorted(icons),
    }


def _color_rows(colors: Iterable[str]) -> list[tuple[str, str]]:
    used: set[str] = set()
    rows: list[tuple[str, str]] = []
    for index, color in enumerate(colors, start=1):
        base = _COLOR_ROLES.get(color, f"color_{index:02d}")
        role = base
        suffix = 2
        while role in used:
            role = f"{base}_{suffix}"
            suffix += 1
        used.add(role)
        rows.append((role, color))
    return rows


def _spec_lock(inventory: dict[str, Any]) -> str:
    """Project a validated flat-project lock from one rendered figure SVG."""

    colors = _color_rows(inventory["colors"])
    sizes = inventory["font_sizes"]
    families = inventory["font_families"]
    color_lines = "\n".join(f"- {role}: {color}" for role, color in colors)
    typography_lines = [
        f"- font_family: {families[0]}",
        f"- body: {_number(sizes[0])}",
        f"- title: {_number(sizes[-1])}",
    ]
    typography_lines.extend(
        f"- scale_{index:02d}: {_number(size)}"
        for index, size in enumerate(sizes[1:-1], start=1)
    )
    typography_lines.extend(
        f"- alternate_{index:02d}_family: {family}"
        for index, family in enumerate(families[1:], start=1)
    )
    icon_inventory = ", ".join(f"tabler-outline/{name}" for name in inventory["icons"])
    return f"""<!-- ppt-master-schema: spec-lock/v1 -->
# Execution Lock

## canvas
- viewBox: 0 0 1280 720
- format: PPT 16:9

## communication
- audience: CCF-A paper reviewers and readers
- objective: Explain one scientific figure accurately in an editable format
- core_message: Native labels and geometry remain deterministic and editable
- consumption_mode: text

## mode
- mode: instructional

## visual_style
- visual_style: soft-rounded

## colors
{color_lines}

## typography
{chr(10).join(typography_lines)}

## icons
- library: tabler-outline
- inventory: {icon_inventory}
- stroke_width: 2

## page_rhythm
- P01: anchor

## pptx_structure
- mode: flat

## forbidden
- Mixing icon libraries
- `mask`, `<style>`, `class`, external CSS, `<foreignObject>`, `textPath`, `@font-face`, `<animate*>`, `<set>`, `<script>` / event attributes, `<iframe>`
- HTML named entities in text; write typography as raw Unicode and escape XML reserved characters
"""


def _design_spec(figure_id: str, final_width_mm: float, inventory: dict[str, Any]) -> str:
    color_rows = "\n".join(
        f"| {role.replace('_', ' ').title()} | {color} | emitted SVG paint |"
        for role, color in _color_rows(inventory["colors"])
    )
    font_rows = "\n".join(
        f"| {role.replace('-', ' ').title()} | {', '.join(_number(size) for size in sizes)} px |"
        for role, sizes in inventory["sizes_by_role"].items()
    )
    icon_rows = "\n".join(
        f"| {name.replace('-', ' ').title()} | templates/icons/tabler-outline/{name}.svg | P01 |"
        for name in inventory["icons"]
    )
    family = inventory["font_families"][0]
    return f"""<!-- ppt-master-schema: design-spec/v1 -->
# {figure_id} - Design Spec

## I. Project Information

| Item | Value |
| --- | --- |
| Project Name | {figure_id} |
| Canvas Format | PPT 16:9 (1280×720) |
| Page Count | 1 |
| Target Audience | CCF-A paper reviewers and readers |
| Communication Intent | Explain a scientific figure accurately |
| Desired Audience Outcome | Understand the figure's evidence and flow |
| Core Message / Ask / Action | Read exact native labels and connections |
| Delivery Context | Paper figure at {final_width_mm:g} mm final width |
| Artifact Afterlife | Editable PowerPoint and vector publication export |
| Reading Mode | text |
| Content Strategy | One deterministic figure per project |
| Design Style | soft-rounded paper figure |
| Created Date | generated by Figure Studio |

## II. Canvas Specification

| Property | Value |
| --- | --- |
| Format | PPT 16:9 |
| Dimensions | 1280×720 |
| viewBox | `0 0 1280 720` |
| Margins | Figure contract controlled |
| Content Area | Figure contract controlled |

## III. Visual Theme

### Theme Style

- **Mode**: instructional
- **Visual style**: soft-rounded
- **Theme**: warm white paper canvas with semantic pastel modules
- **Tone**: compact, precise, publication ready

### Color Scheme

| Role | HEX | Purpose |
| --- | --- | --- |
{color_rows}

## IV. Typography System

### Font Plan

| Role | Chinese | English | Fallback tail |
| --- | --- | --- | --- |
| All emitted text | {family} | {family} | sans-serif |

- All SVG text uses the emitted stack `{family}`.

### Font Size Hierarchy

| Purpose | Size |
| --- | --- |
{font_rows}

## V. Layout Principles

### Page Structure

- **Header area**: contract controlled
- **Content area**: contract controlled
- **Footer area**: omitted unless the contract requires it

### Spacing Specification

| Element | Current Project |
| --- | --- |
| Safe margin | contract controlled |
| Content block gap | 16–24 px |
| Icon-text gap | 8 px when an icon is present |

## VI. Icon Usage Specification

| Purpose | Icon Path | Page |
| --- | --- | --- |
{icon_rows}

## VII. Visualization Reference List

| Page | Template | Path | Summary-quote | Usage |
| --- | --- | --- | --- | --- |

## VIII. Image Resource List

| Filename | Dimensions | Ratio | Purpose | Type | Layout pattern | Acquire Via | Status | Reference | text_policy | page_role |
| --- | --- | --- | --- | --- | --- | --- | --- | --- | --- | --- |

## IX. Content Outline

### Part 1: Figure

#### Slide 01 - {figure_id}

- **Audience move**: Understand the exact scientific labels and connections
- **Layout**: Contract-defined native SVG composition
- **Title**: Contract-defined
- **Core message**: Contract-defined
- **Cover impact**: Native figure geometry + immediate scientific reading path
- **Content**: One editable paper figure

## X. Speaker Notes Requirements

- **Filename**: notes are not required for a paper figure
- **Content**: no speaker notes
"""


def _parse_created_project(output: str) -> Path:
    match = re.search(r"^Project created:\s*(.+)$", output, flags=re.MULTILINE)
    if not match:
        raise ValueError("project_manager.py did not report the created project path")
    return Path(match.group(1).strip()).resolve()


def ensure_project(out_dir: str | Path, figure_id: str, final_width_mm: float) -> Path:
    """Create/refresh a real one-page PPT Master project at ``out_dir``.

    PPT Master's initializer appends a format/date suffix.  We invoke that
    official initializer and both official scaffold commands in a temporary
    sibling, then merge the scaffold into the stable Figure Studio directory.
    """
    project_dir = Path(out_dir).resolve()
    project_dir.mkdir(parents=True, exist_ok=True)
    rendered_svg = project_dir / f"{figure_id}.svg"
    inventory = _svg_inventory(rendered_svg)
    quality = project_dir / "quality"
    quality.mkdir(parents=True, exist_ok=True)

    required_dirs = {
        "svg_output", "svg_final", "images", "icons", "notes", "templates",
        "live_preview", "sources", "analysis", "validation", "exports",
    }
    needs_scaffold = not (project_dir / "README.md").is_file()
    if needs_scaffold:
        init_name = f"argus_{re.sub(r'[^A-Za-z0-9_-]+', '_', figure_id)}_{uuid.uuid4().hex[:8]}"
        completed = _run(
            [PY, PROJECT_MANAGER, "init", init_name, "--format", "ppt169", "--dir", project_dir.parent],
            quality / "project_init.log",
        )
        initialized = _parse_created_project(completed.stdout)
        try:
            _run(
                [PY, PROJECT_MANAGER, "scaffold-spec", initialized],
                quality / "scaffold_spec.log",
            )
            _run(
                [PY, PROJECT_MANAGER, "scaffold-lock", initialized],
                quality / "scaffold_lock.log",
            )
            for child in initialized.iterdir():
                destination = project_dir / child.name
                if child.is_dir():
                    shutil.copytree(child, destination, dirs_exist_ok=True)
                elif child.name not in {"design_spec.md", "spec_lock.md"}:
                    shutil.copy2(child, destination)
        finally:
            shutil.rmtree(initialized, ignore_errors=True)
    for directory in required_dirs:
        (project_dir / directory).mkdir(parents=True, exist_ok=True)

    # Filled only after the versioned scaffolds exist, preserving the required
    # scaffold-first workflow while making repeat builds deterministic.
    (project_dir / "design_spec.md").write_text(
        _design_spec(figure_id, float(final_width_mm), inventory), encoding="utf-8"
    )
    (project_dir / "spec_lock.md").write_text(_spec_lock(inventory), encoding="utf-8")
    return project_dir


def _normalise_checker_report(report: Any) -> dict[str, Any]:
    """Expose stable top-level error/warning lists across checker revisions."""
    if not isinstance(report, dict):
        return {"errors": ["PPT Master returned a non-object JSON report"], "warnings": [], "raw": report}
    errors = list(report.get("errors") or [])
    warnings = list(report.get("warnings") or [])
    if not errors and not warnings:
        for result in report.get("results", report.get("files", [])) or []:
            if isinstance(result, dict):
                errors.extend(result.get("errors") or [])
                warnings.extend(result.get("warnings") or [])
    report["errors"] = errors
    report["warnings"] = warnings
    return report


def check_svg(svg: str | Path) -> dict[str, Any]:
    """Run PPT Master's SVG checker and return its parsed JSON report."""
    svg_path = Path(svg).resolve()
    quality = _quality_dir(svg_path)
    json_path = quality / "pptmaster_check.json"
    completed = _run(
        [PY, SVG_CHECKER, svg_path, "--format", "ppt169", "--json", "--json-output", json_path],
        quality / "pptmaster_check.log",
        check=False,
    )
    try:
        report = json.loads(json_path.read_text(encoding="utf-8"))
    except (OSError, UnicodeError, json.JSONDecodeError) as exc:
        report = {
            "errors": [f"could not parse PPT Master JSON report: {exc}"],
            "warnings": [],
        }
    report = _normalise_checker_report(report)
    report["returncode"] = completed.returncode
    report["log"] = str(quality / "pptmaster_check.log")
    if completed.returncode and not report["errors"]:
        report["errors"].append(
            f"svg_quality_checker.py exited {completed.returncode}; see {quality / 'pptmaster_check.log'}"
        )
    return report


def export_pptx(project_dir: str | Path, out_pptx: str | Path) -> Path:
    """Export one flat native DrawingML presentation from ``svg_output``."""
    project = Path(project_dir).resolve()
    destination = Path(out_pptx).resolve()
    destination.parent.mkdir(parents=True, exist_ok=True)
    _run(
        [PY, SVG_TO_PPTX, project, "-o", destination, "-f", "ppt169", "--pptx-structure", "flat"],
        project / "quality/export_pptx.log",
    )
    if not destination.is_file():
        raise BridgeError(
            "svg_to_pptx.py reported success but did not create the output",
            command=[PY, SVG_TO_PPTX, project],
            log_path=project / "quality/export_pptx.log",
        )
    # PPT Master intentionally omits the SVG's data-pptx-role="background"
    # rectangle from editable slide shapes. Apply the same token to the native
    # slide background so the exported deck still renders the warm paper canvas.
    presentation = Presentation(str(destination))
    for slide in presentation.slides:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xFB, 0xFA, 0xF7)
    presentation.save(str(destination))
    return destination


def pptx_inventory(pptx: str | Path) -> dict[str, Any]:
    presentation = Presentation(str(Path(pptx).resolve()))
    shape_count = 0
    text_frames = 0
    picture_count = 0
    texts: list[str] = []

    def visit(shapes: Any) -> None:
        nonlocal shape_count, text_frames, picture_count
        for shape in shapes:
            shape_count += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
            if getattr(shape, "has_text_frame", False):
                text_frames += 1
                text = str(shape.text)
                if text.strip():
                    texts.append(text)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                visit(shape.shapes)

    for slide in presentation.slides:
        visit(slide.shapes)
    return {
        "shape_count": shape_count,
        "text_frames": text_frames,
        "texts": texts,
        "picture_count": picture_count,
    }


def roundtrip(pptx: str | Path, out_dir: str | Path) -> dict[str, Any]:
    """Convert PPTX back to SVG and inventory editability with python-pptx."""
    pptx_path = Path(pptx).resolve()
    output = Path(out_dir).resolve()
    output.mkdir(parents=True, exist_ok=True)
    quality = _quality_dir(pptx_path)
    _run(
        [PY, PPTX_TO_SVG, pptx_path, "-o", output],
        quality / "roundtrip.log",
    )
    inventory = pptx_inventory(pptx_path)
    result = {
        **inventory,
        "output_dir": str(output),
        "svg_files": [str(path) for path in sorted(output.rglob("*.svg"))],
        "log": str(quality / "roundtrip.log"),
    }
    (quality / "roundtrip_inventory.json").write_text(
        json.dumps(result, indent=2, ensure_ascii=False) + "\n", encoding="utf-8"
    )
    return result


def _parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    commands = parser.add_subparsers(dest="command", required=True)
    ensure = commands.add_parser("ensure-project")
    ensure.add_argument("out_dir", type=Path)
    ensure.add_argument("figure_id")
    ensure.add_argument("final_width_mm", type=float)
    check = commands.add_parser("check-svg")
    check.add_argument("svg", type=Path)
    export = commands.add_parser("export-pptx")
    export.add_argument("project_dir", type=Path)
    export.add_argument("out_pptx", type=Path)
    back = commands.add_parser("roundtrip")
    back.add_argument("pptx", type=Path)
    back.add_argument("out_dir", type=Path)
    return parser


def main(argv: Iterable[str] | None = None) -> int:
    args = _parser().parse_args(argv)
    try:
        if args.command == "ensure-project":
            result: Any = {"project_dir": str(ensure_project(args.out_dir, args.figure_id, args.final_width_mm))}
        elif args.command == "check-svg":
            result = check_svg(args.svg)
        elif args.command == "export-pptx":
            result = {"pptx": str(export_pptx(args.project_dir, args.out_pptx))}
        else:
            result = roundtrip(args.pptx, args.out_dir)
    except (BridgeError, OSError, ValueError) as exc:
        print(f"pptmaster_bridge failed: {exc}", file=sys.stderr)
        return 1
    print(json.dumps(result, indent=2, ensure_ascii=False))
    return 0 if not isinstance(result, dict) or not result.get("errors") else 1


if __name__ == "__main__":
    raise SystemExit(main())
