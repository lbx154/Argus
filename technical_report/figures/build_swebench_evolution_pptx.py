#!/usr/bin/env python3
"""Build the editable PowerPoint figure for the unified SWE-Bench Pro study."""

from __future__ import annotations

import csv
import hashlib
import json
import shutil
import subprocess
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPORT = Path(__file__).resolve().parents[1]
FIGURES = REPORT / "figures"
EVIDENCE = REPORT / "evidence" / "swebench_pro"
SUMMARY_PATH = EVIDENCE / "unified_experiment_summary.json"
WAVES_PATH = EVIDENCE / "argus_wave_efficiency.csv"
PPTX_PATH = FIGURES / "swebench_evolution.pptx"
PDF_PATH = FIGURES / "swebench_evolution.pdf"
SVG_PATH = FIGURES / "swebench_evolution.svg"
PNG_PATH = FIGURES / "swebench_evolution.png"
WINDOWS_PATH = EVIDENCE / "argus_six_wave_windows.csv"
MACROS_PATH = FIGURES / "swebench_metrics.tex"
PROVENANCE_PATH = FIGURES / "swebench_evolution.provenance.json"


WHITE = RGBColor(255, 255, 255)
INK = RGBColor(30, 39, 50)
MUTED = RGBColor(94, 104, 117)
BLUE = RGBColor(49, 91, 206)
DEEP = RGBColor(23, 59, 112)
PALE_BLUE = RGBColor(232, 238, 252)
LIGHT_BLUE = RGBColor(166, 188, 235)
GOLD = RGBColor(195, 138, 32)
PALE_GOLD = RGBColor(252, 246, 229)
GRAY = RGBColor(122, 131, 142)
LIGHT_GRAY = RGBColor(219, 225, 231)
PANEL = RGBColor(248, 250, 252)


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 10,
    color: RGBColor = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font: str = "Aptos",
    margin: float = 0.02,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = LIGHT_GRAY,
    radius: bool = True,
) -> Any:
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    shape.shadow.inherit = False
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor, width: float = 1.2) -> Any:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_circle(slide, cx: float, cy: float, radius: float, color: RGBColor) -> Any:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - radius),
        Inches(cy - radius),
        Inches(2 * radius),
        Inches(2 * radius),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = WHITE
    shape.line.width = Pt(1.0)
    shape.shadow.inherit = False
    return shape


def load_rows() -> list[dict[str, float]]:
    rows: list[dict[str, float]] = []
    with WAVES_PATH.open(encoding="utf-8", newline="") as handle:
        for raw in csv.DictReader(handle):
            rows.append({key: float(value) for key, value in raw.items()})
    return rows


def build_windows(rows: list[dict[str, float]]) -> list[dict[str, Any]]:
    specs = [
        ("W1–6", "Start-up", 1, 6, "normal"),
        ("W7–12", "Early reuse", 7, 12, "normal"),
        ("W13–18", "Composition shift", 13, 18, "normal"),
        ("W19–22", "Mature", 19, 22, "mature"),
        ("W23–24", "Late difficult tasks", 23, 24, "tail"),
    ]
    windows: list[dict[str, Any]] = []
    for label, stage, lower, upper, kind in specs:
        selected = [row for row in rows if lower <= int(row["wave"]) <= upper]
        tasks = sum(int(row["completed"]) for row in selected)
        weighted = lambda key: sum(row[key] * row["completed"] for row in selected) / tasks
        windows.append(
            {
                "label": label,
                "stage": stage,
                "kind": kind,
                "waves": ",".join(str(int(row["wave"])) for row in selected),
                "tasks": tasks,
                "solve_input_tokens_mean": weighted("solve_input_tokens_mean"),
                "active_seconds_mean": weighted("solve_agent_seconds_mean"),
            }
        )
    start = windows[0]
    for window in windows:
        window["token_index"] = 100.0 * window["solve_input_tokens_mean"] / start["solve_input_tokens_mean"]
        window["time_index"] = 100.0 * window["active_seconds_mean"] / start["active_seconds_mean"]
    return windows


def write_windows(rows: list[dict[str, Any]]) -> None:
    fields = [
        "label",
        "stage",
        "kind",
        "waves",
        "tasks",
        "solve_input_tokens_mean",
        "active_seconds_mean",
        "token_index",
        "time_index",
    ]
    with WINDOWS_PATH.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=fields)
        writer.writeheader()
        writer.writerows(rows)


def draw_accuracy_panel(slide, summary: dict[str, Any]) -> None:
    x, y, w, h = 0.35, 1.15, 3.2, 5.1
    add_rect(slide, x, y, w, h, fill=PANEL)
    add_text(slide, "(a) Full-suite comparison", x + 0.18, y + 0.10, 2.8, 0.32, size=12, bold=True, color=DEEP)
    aggregate = summary["aggregate_comparison"]
    direct = 100.0 * aggregate["direct_copilot_accuracy_approx"]
    argus = 100.0 * aggregate["argus_accuracy_approx"]
    add_text(slide, "Accuracy on 731 tasks", x + 0.18, y + 0.53, 2.6, 0.25, size=9.5, bold=True)
    bar_x, bar_w = x + 0.18, 2.55
    for label, value, yy, color in (
        ("Direct Copilot", direct, y + 0.90, GRAY),
        ("Argus", argus, y + 1.55, BLUE),
    ):
        add_text(slide, label, bar_x, yy - 0.27, 1.4, 0.22, size=8.7, bold=True if label == "Argus" else False)
        add_rect(slide, bar_x, yy, bar_w, 0.28, fill=WHITE, line=LIGHT_GRAY, radius=False)
        add_rect(slide, bar_x, yy, bar_w * value / 100.0, 0.28, fill=color, line=color, radius=False)
        add_text(slide, f"≈{value:.0f}%", bar_x + bar_w * value / 100.0 - 0.62, yy - 0.01, 0.56, 0.28, size=10, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    add_text(slide, "+19 percentage points", x + 0.18, y + 2.03, 2.55, 0.33, size=13.5, bold=True, color=BLUE)

    add_line(slide, x + 0.18, y + 2.48, x + w - 0.18, y + 2.48, LIGHT_GRAY, 0.8)
    add_text(slide, "Normalized aggregate tokens", x + 0.18, y + 2.65, 2.6, 0.25, size=9.5, bold=True)
    token_ratio = float(aggregate["argus_to_direct_total_token_ratio_approx"])
    scale = 1.5
    for label, value, yy, color in (
        ("Direct Copilot", 1.0, y + 3.03, GRAY),
        ("Argus", token_ratio, y + 3.60, BLUE),
    ):
        add_text(slide, label, bar_x, yy - 0.24, 1.4, 0.20, size=8.5)
        add_rect(slide, bar_x, yy, bar_w, 0.22, fill=WHITE, line=LIGHT_GRAY, radius=False)
        add_rect(slide, bar_x, yy, bar_w * value / scale, 0.22, fill=color, line=color, radius=False)
        add_text(slide, f"{value:.2f}×", bar_x + bar_w + 0.03, yy - 0.03, 0.45, 0.28, size=9.5, bold=True, color=color)
    add_text(
        slide,
        "Only the aggregate Token ratio is available;\nCopilot per-wave Token/time was not logged.",
        x + 0.18,
        y + 4.13,
        2.78,
        0.55,
        size=8.1,
        color=MUTED,
    )


def draw_series_panel(
    slide,
    windows: list[dict[str, Any]],
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    value_key: str,
    unit: str,
    y_max: float,
    grid_step: float,
    reduction: float,
    best_reduction: float | None = None,
) -> None:
    add_rect(slide, x, y, w, h, fill=WHITE)
    add_text(slide, title, x + 0.16, y + 0.08, w - 0.32, 0.30, size=11.5, bold=True, color=DEEP)
    plot_x, plot_y = x + 0.56, y + 0.48
    plot_w, plot_h = w - 1.10, h - 0.94
    tail_x = plot_x + plot_w * 4 / 4
    add_rect(slide, tail_x - 0.35, plot_y, 0.70, plot_h, fill=PALE_GOLD, line=PALE_GOLD, radius=False)
    for tick in [grid_step * i for i in range(int(y_max / grid_step) + 1)]:
        py = plot_y + plot_h - plot_h * tick / y_max
        add_line(slide, plot_x, py, plot_x + plot_w, py, LIGHT_GRAY, 0.65)
        tick_text = f"{tick / 1_000_000:g}" if value_key == "solve_input_tokens_mean" else f"{tick / 60:g}"
        add_text(slide, tick_text, x + 0.08, py - 0.11, 0.38, 0.22, size=7.5, color=MUTED, align=PP_ALIGN.RIGHT)
    points: list[tuple[float, float]] = []
    for index, window in enumerate(windows):
        px = plot_x + plot_w * index / (len(windows) - 1)
        value = float(window[value_key])
        py = plot_y + plot_h - plot_h * value / y_max
        points.append((px, py))
        if index:
            color = GOLD if window["kind"] == "tail" else LIGHT_BLUE
            add_line(slide, points[index - 1][0], points[index - 1][1], px, py, color, 2.0)
        color = GOLD if window["kind"] == "tail" else BLUE
        add_circle(slide, px, py, 0.075, color)
        value_text = f"{value / 1_000_000:.2f}M" if value_key == "solve_input_tokens_mean" else f"{value / 60:.2f}m"
        add_text(slide, value_text, px - 0.35, py - 0.32, 0.70, 0.22, size=8.0, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, window["label"], px - 0.43, plot_y + plot_h + 0.07, 0.86, 0.20, size=7.6, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_text(slide, window["stage"], px - 0.55, plot_y + plot_h + 0.27, 1.10, 0.20, size=7.0, color=MUTED, align=PP_ALIGN.CENTER)
    mature = points[3]
    add_text(
        slide,
        f"Start-up → mature: −{reduction:.0f}%",
        x + w - 2.03,
        y + 0.10,
        1.82,
        0.27,
        size=9.1,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.RIGHT,
    )
    if best_reduction is not None:
        add_text(slide, f"Best six-wave regime: −{best_reduction:.0f}%", x + 0.55, y + 0.37, 2.1, 0.22, size=7.8, color=MUTED)
    if value_key == "active_seconds_mean":
        spike = points[2]
        add_text(slide, "composition shift", spike[0] - 0.57, spike[1] + 0.12, 1.14, 0.20, size=7.2, color=GOLD, bold=True, align=PP_ALIGN.CENTER)


def write_macros(summary: dict[str, Any], windows: list[dict[str, Any]]) -> None:
    aggregate = summary["aggregate_comparison"]
    start, best, mature = windows[0], windows[2], windows[3]
    token_reduction = 100.0 * (1.0 - mature["solve_input_tokens_mean"] / start["solve_input_tokens_mean"])
    time_reduction = 100.0 * (1.0 - mature["active_seconds_mean"] / start["active_seconds_mean"])
    best_token_reduction = 100.0 * (1.0 - best["solve_input_tokens_mean"] / start["solve_input_tokens_mean"])
    values = {
        "SWEProTasks": int(summary["tasks"]),
        "SWEProDirectAccuracy": f"{100 * aggregate['direct_copilot_accuracy_approx']:.0f}",
        "SWEProArgusAccuracy": f"{100 * aggregate['argus_accuracy_approx']:.0f}",
        "SWEProAccuracyDelta": f"{aggregate['accuracy_delta_percentage_points_approx']:.0f}",
        "SWEProTokenRatio": f"{aggregate['argus_to_direct_total_token_ratio_approx']:.2f}",
        "SWEProStartupTokensM": f"{start['solve_input_tokens_mean'] / 1_000_000:.2f}",
        "SWEProMatureTokensM": f"{mature['solve_input_tokens_mean'] / 1_000_000:.2f}",
        "SWEProTokenReduction": f"{token_reduction:.0f}",
        "SWEProBestTokenReduction": f"{best_token_reduction:.0f}",
        "SWEProStartupMinutes": f"{start['active_seconds_mean'] / 60:.2f}",
        "SWEProMatureMinutes": f"{mature['active_seconds_mean'] / 60:.2f}",
        "SWEProTimeReduction": f"{time_reduction:.0f}",
        "SWEProCompletedWaves": len(load_rows()),
    }
    MACROS_PATH.write_text(
        "".join(f"\\newcommand{{\\{key}}}{{{value}}}\n" for key, value in values.items()),
        encoding="utf-8",
    )


def main() -> int:
    summary = json.loads(SUMMARY_PATH.read_text(encoding="utf-8"))
    wave_rows = load_rows()
    windows = build_windows(wave_rows)
    write_windows(windows)
    write_macros(summary, windows)

    start, best, mature = windows[0], windows[2], windows[3]
    token_reduction = 100.0 * (1.0 - mature["solve_input_tokens_mean"] / start["solve_input_tokens_mean"])
    time_reduction = 100.0 * (1.0 - mature["active_seconds_mean"] / start["active_seconds_mean"])
    best_token_reduction = 100.0 * (1.0 - best["solve_input_tokens_mean"] / start["solve_input_tokens_mean"])

    prs = Presentation()
    prs.slide_width = Inches(12)
    prs.slide_height = Inches(6.75)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_rect(slide, 0, 0, 12, 0.065, fill=BLUE, line=BLUE, radius=False)
    add_text(slide, "SWE-Bench Pro: accuracy gain and runtime convergence", 0.35, 0.20, 11.2, 0.42, size=20, bold=True, color=INK)
    add_text(
        slide,
        "One continuous 731-task evaluation · GPT-5.5/xhigh · Skill and Wiki state updated throughout the run",
        0.37,
        0.64,
        11.1,
        0.28,
        size=9.5,
        color=MUTED,
    )
    draw_accuracy_panel(slide, summary)
    draw_series_panel(
        slide,
        windows,
        x=3.75,
        y=1.15,
        w=7.90,
        h=2.40,
        title="(b) Argus solve input per task",
        value_key="solve_input_tokens_mean",
        unit="M tokens",
        y_max=4_200_000,
        grid_step=1_000_000,
        reduction=token_reduction,
        best_reduction=best_token_reduction,
    )
    draw_series_panel(
        slide,
        windows,
        x=3.75,
        y=3.80,
        w=7.90,
        h=2.40,
        title="(c) Argus active workflow time per task",
        value_key="active_seconds_mean",
        unit="minutes",
        y_max=720,
        grid_step=120,
        reduction=time_reduction,
    )
    add_text(
        slide,
        "Task-weighted window means. Two incomplete Waves are omitted; W23–24 are shown separately as late difficult-task stress. Active time excludes orchestration wait, environment preparation, external verification, infrastructure recovery, and post-task knowledge maintenance.",
        0.38,
        6.40,
        11.2,
        0.22,
        size=7.2,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    prs.save(PPTX_PATH)

    libreoffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not libreoffice:
        raise SystemExit("LibreOffice is required to export the PPTX to vector PDF")
    if PDF_PATH.exists():
        PDF_PATH.unlink()
    result = subprocess.run(
        [libreoffice, "--headless", "--convert-to", "pdf", "--outdir", str(FIGURES), str(PPTX_PATH)],
        capture_output=True,
        text=True,
        check=False,
    )
    if result.returncode != 0 or not PDF_PATH.is_file():
        raise SystemExit(f"PowerPoint PDF export failed: {result.stderr}")
    ghostscript = shutil.which("gs")
    if ghostscript:
        compatible = FIGURES / "swebench_evolution.compat.pdf"
        subprocess.run(
            [
                ghostscript,
                "-q",
                "-dNOPAUSE",
                "-dBATCH",
                "-sDEVICE=pdfwrite",
                "-dCompatibilityLevel=1.5",
                f"-sOutputFile={compatible}",
                str(PDF_PATH),
            ],
            check=True,
        )
        compatible.replace(PDF_PATH)
    subprocess.run(["pdftocairo", "-svg", str(PDF_PATH), str(SVG_PATH)], check=True)
    preview_prefix = PNG_PATH.with_suffix("")
    subprocess.run(["pdftoppm", "-singlefile", "-png", "-r", "180", str(PDF_PATH), str(preview_prefix)], check=True)

    provenance = {
        "figure_id": "swebench-unified-evolution",
        "reader_question": "Does persistent Skill/Wiki evolution improve full-suite accuracy and reduce mature-run Argus resource use?",
        "claim": f"Argus reaches approximately 78% versus 59% for Direct Copilot at 1.41x aggregate tokens; W19-22 uses {token_reduction:.0f}% fewer solve tokens and {time_reduction:.0f}% less active time than W1-6.",
        "scope": "One 731-task experiment. Copilot per-wave resource traces are unavailable; Argus windows are observational and W23-24 are shown separately as late difficult-task stress.",
        "inputs": {SUMMARY_PATH.name: sha256(SUMMARY_PATH), WAVES_PATH.name: sha256(WAVES_PATH)},
        "outputs": {path.name: sha256(path) for path in (PPTX_PATH, PDF_PATH, SVG_PATH, PNG_PATH, WINDOWS_PATH, MACROS_PATH)},
        "editable_source": PPTX_PATH.name,
    }
    PROVENANCE_PATH.write_text(json.dumps(provenance, indent=2, sort_keys=True) + "\n", encoding="utf-8")
    print(PPTX_PATH)
    print(PDF_PATH)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
